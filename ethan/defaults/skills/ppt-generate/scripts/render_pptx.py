#!/usr/bin/env python3
"""render_pptx.py — 瘦身版 PPTist JSON → 原生可编辑 pptx 渲染器。

用法:
    python3 render_pptx.py deck.json -o out.pptx [--theme NAME] [--check]
    python3 render_pptx.py <项目目录> -o out.pptx   # 含 deck.json + pages/*.json，按文件名排序合并

- deck.json 格式见 references/schema.md（画布默认 1000x562.5 px，1px = 12192 EMU，字号 pt = px * 0.96）
- 项目目录结构见 project_loader.py docstring（逐页生成工作流：deck.json 元信息 + 每页一个 Slide JSON）
- --check 只校验不渲染，exit 0 = 通过
- 依赖 python-pptx；缺失时自动 `pip install --user python-pptx` 后重试
"""
from __future__ import annotations

import argparse
import copy
import hashlib
import json
import math
import os
import re
import subprocess
import sys
import tempfile
import urllib.request
from pathlib import Path

from project_loader import load_deck

# ---------------------------------------------------------------------------
# 依赖自举：python-pptx（ Pillow 随其自动安装 ）
# ---------------------------------------------------------------------------

def _ensure_pptx():
    missing = []
    for mod in ("pptx", "latex2mathml", "mathml2omml"):
        try:
            __import__(mod)
        except ImportError:
            missing.append(mod)
    if not missing:
        return
    pkg_map = {"pptx": "python-pptx", "latex2mathml": "latex2mathml", "mathml2omml": "mathml2omml"}
    pkgs = [pkg_map[m] for m in missing]
    print(f"[render_pptx] 缺少依赖 {pkgs}，尝试 pip 安装 ...", file=sys.stderr)
    # venv 内直接装（--user 会报错）；非 venv 用 --user 避免污染系统环境
    in_venv = sys.prefix != getattr(sys, "base_prefix", sys.prefix)
    cmd = [sys.executable, "-m", "pip", "install", "--quiet", *([] if in_venv else ["--user"]), *pkgs]
    try:
        subprocess.check_call(cmd)
        import site

        user_site = site.getusersitepackages()
        if user_site not in sys.path:
            sys.path.insert(0, user_site)
        import pptx  # noqa: F401
        print("[render_pptx] 依赖安装成功", file=sys.stderr)
    except Exception as e:  # noqa: BLE001
        print(
            "[render_pptx] 自动安装失败: %s\n请手动执行: pip3 install python-pptx latex2mathml mathml2omml" % e,
            file=sys.stderr,
        )
        sys.exit(2)


_ensure_pptx()

from pptx import Presentation  # noqa: E402
from pptx.chart.data import CategoryChartData, XyChartData  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION  # noqa: E402
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE  # noqa: E402
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402
from pptx.util import Emu, Pt  # noqa: E402

try:
    from PIL import Image as _PILImage
except ImportError:  # pragma: no cover - Pillow 是 python-pptx 的依赖，正常不会缺
    _PILImage = None

# ---------------------------------------------------------------------------
# 常量
# ---------------------------------------------------------------------------

DEFAULT_CANVAS_W = 1000.0
DEFAULT_CANVAS_H = 562.5
SLIDE_W_EMU = 12192000  # 13.333 in，16:9
PT_PER_PX = 0.96  # 1000px = 13.333in → 1px = 0.96pt

# 上下标基线偏移（%）：与前端 slide.tsx 的 0.7 缩放配套（预览下标略深、上标略浅是 pptx 惯例）
SUB_BASELINE_PCT = -25
SUP_BASELINE_PCT = 30

SLIDE_TYPES = {"cover", "contents", "transition", "content", "end"}
TEXT_TYPES = {
    "title", "subtitle", "content", "item", "itemTitle",
    "notes", "header", "footer", "partNumber", "itemNumber",
}
IMAGE_TYPES = {"pageFigure", "itemFigure", "background"}

# OOXML 预设几何（a:prstGeom prst）白名单 —— 与 references/schema.md 表格一致
PRESET_SHAPES = {
    # 矩形
    "rect", "roundRect", "round1Rect", "round2SameRect", "round2DiagRect",
    "snip1Rect", "snip2SameRect", "snip2DiagRect", "snipRoundRect",
    # 基础
    "ellipse", "triangle", "rtTriangle", "diamond", "parallelogram", "trapezoid",
    "pentagon", "hexagon", "heptagon", "octagon", "plus", "donut",
    # 箭头
    "rightArrow", "leftArrow", "upArrow", "downArrow", "leftRightArrow",
    "upDownArrow", "bentArrow", "chevron", "notchedRightArrow",
    # 星/旗
    "star4", "star5", "star6", "star8", "ribbon", "ribbon2", "wave", "doubleWave",
    # 对话/标注
    "wedgeRectCallout", "wedgeRoundRectCallout", "wedgeEllipseCallout", "cloudCallout",
    # 其他
    "heart", "cloud", "sun", "moon", "cube", "can", "teardrop", "frame",
    "halfFrame", "corner", "diagStripe", "foldedCorner", "smileyFace",
    "lightningBolt", "bracketPair", "bracePair", "blockArc", "pie", "chord",
}

CHART_TYPES = {
    "column", "bar", "line", "pie", "ring", "area", "radar", "scatter",
}

ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}
VALIGN_MAP = {
    "top": MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}
DASH_MAP = {"solid": None, "dashed": "dash", "dotted": "sysDot"}

DEFAULT_THEME = {
    "name": "default",
    "backgroundColor": "#FFFFFF",
    "themeColors": ["#1E40AF", "#3B82F6", "#93C5FD", "#F59E0B", "#10B981"],
    "fontColor": "#1F2937",
    "fontName": "Microsoft YaHei",
    "outline": {"style": "solid", "width": 1, "color": "#D1D5DB"},
    "shadow": {"h": 0, "v": 2, "blur": 8, "color": "#00000014"},
    "typography": {
        "title": {"fontSize": 28, "color": "#111827", "bold": True},
        "subtitle": {"fontSize": 16, "color": "#4B5563"},
        "content": {"fontSize": 14, "color": "#1F2937"},
        "item": {"fontSize": 14, "color": "#1F2937"},
        "itemTitle": {"fontSize": 16, "color": "#111827", "bold": True},
        "notes": {"fontSize": 10, "color": "#9CA3AF"},
        "header": {"fontSize": 10, "color": "#6B7280"},
        "footer": {"fontSize": 10, "color": "#6B7280"},
        "partNumber": {"fontSize": 60, "color": "#1E40AF", "bold": True},
        "itemNumber": {"fontSize": 16, "color": "#1E40AF", "bold": True},
    },
}

_TMP_IMG_DIR = Path(tempfile.gettempdir()) / "ppt_render_imgs"

# OMML（Office Math Markup Language）命名空间 —— latex 元素用
M_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
# DrawingML 2010 扩展：pptx 文本体内的数学区必须包在 <a14:m> 里（对照 pandoc 输出）
A14_NS = "http://schemas.microsoft.com/office/drawing/2010/main"
# 标记兼容性：a14:m 必须再包一层 mc:AlternateContent（Choice/Fallback），
# 否则 Office 365 不渲染公式、Keynote 直接拒收整个文件


# ---------------------------------------------------------------------------
# 基础工具
# ---------------------------------------------------------------------------

class DeckError(Exception):
    """deck JSON 校验/渲染错误"""


def parse_color(value: str):
    """#RGB / #RRGGBB / #RRGGBBAA → (RGBColor, alpha_float|None)。alpha 为 CSS 8 位 hex 的末两位。"""
    if not isinstance(value, str):
        raise DeckError(f"颜色必须是字符串: {value!r}")
    s = value.strip().lstrip("#")
    if len(s) == 3:
        s = "".join(c * 2 for c in s)
    if len(s) == 6:
        return RGBColor.from_string(s.upper()), None
    if len(s) == 8:
        rgb, a = s[:6], s[6:]
        return RGBColor.from_string(rgb.upper()), int(a, 16) / 255.0
    raise DeckError(f"无法解析颜色: {value!r}")


def rel_luminance(rgb: RGBColor) -> float:
    """相对亮度（0-1），用于判断深/浅色背景。"""
    r, g, b = (int(v) / 255 for v in rgb)
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def shift_color(rgb: RGBColor, amount: float) -> RGBColor:
    """整体提亮（amount>0）或压暗（amount<0），amount 为 0-1 的幅度。"""
    def f(v):
        return max(0, min(255, round(int(v) + 255 * amount)))
    return RGBColor(f(rgb[0]), f(rgb[1]), f(rgb[2]))


def contrast_text_color(bg: RGBColor, dark="#1F2937", light="#FFFFFF") -> str:
    """根据底色亮度选黑/白文字。"""
    return dark if rel_luminance(bg) > 0.5 else light


def px_to_emu(px: float, emu_per_px: float) -> Emu:
    return Emu(int(round(px * emu_per_px)))

def px_to_pt(px: float) -> Pt:
    return Pt(px * PT_PER_PX)


def _sub(parent, tag: str, **attrs):
    el = parent.makeelement(qn(tag), {})
    for k, v in attrs.items():
        el.set(k, str(v))
    parent.append(el)
    return el


def _set_solid_fill_alpha(fill_parent_el, alpha: float):
    """给 a:srgbClr 追加 alpha（0~1 → 0~100000）。"""
    srgb = fill_parent_el.find(qn("a:srgbClr"))
    if srgb is not None and alpha is not None:
        for old in srgb.findall(qn("a:alpha")):
            srgb.remove(old)
        _sub(srgb, "a:alpha", val=str(int(alpha * 100000)))


def apply_solid_fill(fill_owner, color_value, alpha_extra=None):
    """fill_owner: FillFormat。解析 #RRGGBB[AA] 并写入透明度。"""
    rgb, alpha = parse_color(color_value)
    if alpha_extra is not None:
        alpha = alpha_extra if alpha is None else alpha * alpha_extra
    fill_owner.solid()
    fill_owner.fore_color.rgb = rgb
    if alpha is not None:
        srgb = fill_owner._xPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        _set_solid_fill_alpha(srgb.getparent() if srgb is None else srgb.getparent(), alpha)


def apply_gradient_fill(fill_owner, gradient: dict, emu_per_px: float):
    """线性渐变走 python-pptx API；径向渐变改 XML。"""
    colors = gradient.get("colors") or []
    if not colors:
        raise DeckError("gradient.colors 不能为空")
    fill_owner.gradient()
    stops = fill_owner.gradient_stops
    # python-pptx 默认给 2 个 stop；数量不一致时改 XML
    if len(colors) != len(stops):
        grad = fill_owner._xPr.find(qn("a:gradFill"))
        gs_lst = grad.find(qn("a:gsLst"))
        for gs in list(gs_lst):
            gs_lst.remove(gs)
        for c in colors:
            rgb, alpha = parse_color(c["color"])
            gs = _sub(gs_lst, "a:gs", pos=str(int(c.get("pos", 0) * 1000)))
            clr = _sub(gs, "a:srgbClr", val=str(rgb))
            if alpha is not None:
                _sub(clr, "a:alpha", val=str(int(alpha * 100000)))
    else:
        for stop, c in zip(stops, colors):
            rgb, alpha = parse_color(c["color"])
            stop.position = max(0.0, min(1.0, c.get("pos", 0) / 100.0))
            stop.color.rgb = rgb
            if alpha is not None:
                srgb = stop._gs.find(qn("a:srgbClr"))
                _set_solid_fill_alpha(srgb.getparent(), alpha)
    if gradient.get("type", "linear") == "radial":
        # python-pptx 只支持线性角度；径向改 XML：a:lin → a:path
        grad = fill_owner._xPr.find(qn("a:gradFill"))
        lin = grad.find(qn("a:lin"))
        if lin is not None:
            grad.remove(lin)
        path = _sub(grad, "a:path", path="circle")
        _sub(path, "a:fillToRect", l="50000", t="50000", r="50000", b="50000")
    else:
        try:
            fill_owner.gradient_angle = float(gradient.get("rotate", 0))
        except Exception:  # noqa: BLE001 - 部分版本只读，忽略
            pass


def apply_shadow(sp_pr, shadow: dict, emu_per_px: float):
    """注入 a:outerShdw。shadow: {h, v, blur(px), color:#RRGGBB[AA]}"""
    h, v = float(shadow.get("h", 0)), float(shadow.get("v", 0))
    blur = float(shadow.get("blur", 0))
    rgb, alpha = parse_color(shadow.get("color", "#00000033"))
    dist = math.hypot(h, v)
    direction = int((math.degrees(math.atan2(v, h)) % 360) * 60000)
    old = sp_pr.find(qn("a:effectLst"))
    if old is not None:
        sp_pr.remove(old)
    effect = _sub(sp_pr, "a:effectLst")
    shdw = _sub(
        effect, "a:outerShdw",
        blurRad=str(int(blur * emu_per_px)),
        dist=str(int(dist * emu_per_px)),
        dir=str(direction), rotWithShape="0",
    )
    clr = _sub(shdw, "a:srgbClr", val=str(rgb))
    _sub(clr, "a:alpha", val=str(int((alpha if alpha is not None else 1.0) * 100000)))


def apply_line_format(line_owner, spec: dict, emu_per_px: float, default=None):
    """设置线颜色/宽度/虚线。spec: {style, width(px), color}；spec 为空时用 default 或隐藏线条。"""
    from pptx.oxml.ns import qn as _qn

    spec = spec if spec is not None else default
    ln = line_owner  # LineFormat
    if not spec:
        ln.fill.background()
        return
    rgb, alpha = parse_color(spec.get("color", "#D1D5DB"))
    ln.color.rgb = rgb
    ln.width = Emu(int(spec.get("width", 1) * emu_per_px))
    ln_el = ln._get_or_add_ln()
    if alpha is not None:
        solid = ln_el.find(_qn("a:solidFill"))
        if solid is not None:
            _set_solid_fill_alpha(solid, alpha)
    dash = DASH_MAP.get(spec.get("style", "solid"))
    # 清理旧的 prstDash 再加
    for old in ln_el.findall(_qn("a:prstDash")):
        ln_el.remove(old)
    if dash:
        _sub(ln_el, "a:prstDash", val=dash)


def set_run_font(run, font_name: str, latin_name: str | None = None):
    """设置字体：ea/cs 用 font_name（中文），latin 用 latin_name（西文，缺省同 font_name）。"""
    run.font.name = latin_name or font_name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = _sub(rPr, tag)
        e.set("typeface", font_name)


def set_flip(shape, flip_h=False, flip_v=False):
    """水平/垂直翻转（python-pptx 未暴露，改 xfrm 属性）。"""
    if not (flip_h or flip_v):
        return
    xfrm = shape._element.spPr.find(qn("a:xfrm"))
    if xfrm is None:
        return
    if flip_h:
        xfrm.set("flipH", "1")
    if flip_v:
        xfrm.set("flipV", "1")


def set_picture_opacity(pic, opacity: float):
    if opacity >= 0.999:
        return
    blip = pic._element.find(qn("p:blipFill") + "/" + qn("a:blip"))
    if blip is None:
        return
    _sub(blip, "a:alphaModFix", amt=str(int(opacity * 100000)))


def set_picture_round(pic, radius_px: float, w_px: float, h_px: float):
    """圆角图片：custGeom → prstGeom roundRect。"""
    if radius_px <= 0:
        return
    sp_pr = pic._element.spPr
    for tag in ("a:custGeom", "a:prstGeom"):
        e = sp_pr.find(qn(tag))
        if e is not None:
            sp_pr.remove(e)
    adj = int(min(radius_px / max(1.0, min(w_px, h_px)), 0.5) * 100000)
    geom = _sub(sp_pr, "a:prstGeom", prst="roundRect")
    av = _sub(geom, "a:avLst")
    _sub(av, "a:gd", name="adj", fmla=f"val {adj}")


def set_shape_geometry(shape, prst_name: str, adjusts=None):
    """把 autoshape 的 prstGeom 换成目标预设形状。"""
    sp_pr = shape._element.spPr
    geom = sp_pr.find(qn("a:prstGeom"))
    if geom is None:
        geom = _sub(sp_pr, "a:prstGeom")
    geom.set("prst", prst_name)
    for old in geom.findall(qn("a:avLst")):
        geom.remove(old)
    av = _sub(geom, "a:avLst")
    for i, val in enumerate(adjusts or [], start=1):
        _sub(av, "a:gd", name=f"adj{i}", fmla=f"val {int(val)}")


def set_vertical_text(text_frame):
    bodyPr = text_frame._txBody.find(qn("a:bodyPr"))
    bodyPr.set("vert", "eaVert")


def set_bullet(paragraph, kind: str):
    """kind: 'bullet' | 'number'"""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", "228600")
    pPr.set("indent", "-228600")
    if kind == "number":
        _sub(pPr, "a:buAutoNum", type="arabicPeriod")
    else:
        _sub(pPr, "a:buFont", typeface="Arial")
        _sub(pPr, "a:buChar", char="•")


def set_arrowheads(connector, points):
    """points: [start, end]，'' / 'arrow' / 'dot'。"""
    ln = connector.line._get_or_add_ln()
    type_map = {"arrow": "triangle", "dot": "oval"}
    for tag, kind in (("a:headEnd", points[0] if points else ""), ("a:tailEnd", points[1] if points and len(points) > 1 else "")):
        for old in ln.findall(qn(tag)):
            ln.remove(old)
        if kind in type_map:
            _sub(ln, tag, type=type_map[kind], w="med", len="med")


def set_strikethrough(run):
    rPr = run._r.get_or_add_rPr()
    rPr.set("strike", "sngStrike")


def set_baseline(run, pct: int):
    """上/下标：rPr baseline（千分比），正值上标、负值下标。"""
    rPr = run._r.get_or_add_rPr()
    rPr.set("baseline", str(pct * 1000))


def set_word_space(run, px: float):
    """字间距 px → rPr spc（1/100 pt）。"""
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(px * PT_PER_PX * 100)))


# ---------------------------------------------------------------------------
# 主题
# ---------------------------------------------------------------------------

def load_theme(deck: dict, cli_theme: str | None, script_dir: Path) -> dict:
    theme_spec = cli_theme or deck.get("theme")
    theme = copy.deepcopy(DEFAULT_THEME)
    loaded = None
    if isinstance(theme_spec, dict):
        loaded = theme_spec
    elif isinstance(theme_spec, str):
        path = script_dir / "themes" / f"{theme_spec}.json"
        if not path.is_file():
            raise DeckError(f"主题不存在: {theme_spec}（{path}）")
        loaded = json.loads(path.read_text(encoding="utf-8"))
    if loaded:
        for k, v in loaded.items():
            if k == "typography" and isinstance(v, dict):
                theme["typography"].update(v)
            else:
                theme[k] = v
    return theme


def typo_default(theme: dict, text_type: str | None) -> dict:
    typo = theme.get("typography") or {}
    base = dict(typo.get("content") or {})
    if text_type and text_type in typo:
        base.update(typo[text_type])
    base.setdefault("fontSize", 14)
    base.setdefault("color", theme.get("fontColor", "#1F2937"))
    base.setdefault("bold", False)
    return base


# ---------------------------------------------------------------------------
# 元素渲染
# ---------------------------------------------------------------------------

def _pick(d: dict, *keys, default=None):
    """多键别名取值（对齐 PPTist 的 fontsize/fontname 等小写命名）。"""
    for k in keys:
        v = d.get(k)
        if v is not None:
            return v
    return default


def render_paragraphs(text_frame, paragraphs, theme, text_type=None, el_defaults=None):
    defaults = typo_default(theme, text_type)
    if el_defaults:
        defaults.update(el_defaults)
    font_name = defaults.get("fontName") or theme.get("fontName", "Microsoft YaHei")
    latin_name = defaults.get("latinFontName") or theme.get("latinFontName")
    for i, para in enumerate(paragraphs):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.alignment = ALIGN_MAP.get(para.get("align", "left"), PP_ALIGN.LEFT)
        if para.get("lineHeight"):
            p.line_spacing = float(para["lineHeight"])
        if para.get("spaceBefore") is not None:
            p.space_before = px_to_pt(float(para.get("spaceBefore") or 0))
        space_after = _pick(para, "spaceAfter", "paragraphSpace")
        if space_after is not None:
            p.space_after = px_to_pt(float(space_after or 0))
        bullet = para.get("bullet")
        if bullet in ("bullet", "number"):
            set_bullet(p, bullet)
        elif bullet is True:
            set_bullet(p, "bullet")
        runs = para.get("runs") or []
        if not runs and para.get("text"):
            runs = [{"text": para["text"]}]
        for r in runs:
            run = p.add_run()
            run.text = str(r.get("text", ""))
            f = run.font
            f.size = px_to_pt(float(_pick(r, "fontSize", "fontsize") or defaults["fontSize"]))
            f.bold = bool(r.get("bold", defaults.get("bold", False)))
            f.italic = bool(r.get("italic", defaults.get("italic", False)))
            if r.get("underline"):
                f.underline = True
            rgb, alpha = parse_color(r.get("color") or defaults["color"])
            f.color.rgb = rgb
            run_font = _pick(r, "fontName", "fontname")
            # run 显式指定 fontName 时中西文都用它；否则中文用 font_name、西文用主题 latinFontName
            set_run_font(run, run_font or font_name, None if run_font else latin_name)
            if r.get("strikethrough"):
                set_strikethrough(run)
            if r.get("sub") or r.get("sup"):
                set_baseline(run, SUB_BASELINE_PCT if r.get("sub") else SUP_BASELINE_PCT)
            if r.get("wordSpace"):
                set_word_space(run, float(r["wordSpace"]))


def setup_text_frame(text_frame, el, theme, emu_per_px, text_type=None):
    inset = el.get("inset") or [10, 10, 10, 10]
    text_frame.margin_top = px_to_emu(inset[0], emu_per_px)
    text_frame.margin_right = px_to_emu(inset[1], emu_per_px)
    text_frame.margin_bottom = px_to_emu(inset[2], emu_per_px)
    text_frame.margin_left = px_to_emu(inset[3], emu_per_px)
    text_frame.word_wrap = True
    # 垂直对齐：元素级 vAlign 优先；形状内嵌文本走 text.align（schema 约定）；
    # render_shape 直接把 text_spec 传进来时读它自身的 align
    v_align = el.get("vAlign") or (el.get("text") or {}).get("align") or el.get("align") or "top"
    text_frame.vertical_anchor = VALIGN_MAP.get(v_align, MSO_ANCHOR.TOP)
    if el.get("vertical"):
        set_vertical_text(text_frame)


def render_text(slide, el, theme, emu_per_px):
    box = slide.shapes.add_textbox(
        px_to_emu(el["left"], emu_per_px), px_to_emu(el["top"], emu_per_px),
        px_to_emu(el["width"], emu_per_px), px_to_emu(el["height"], emu_per_px),
    )
    if el.get("name") or el.get("id"):
        box.name = el.get("name") or el["id"]
    box.rotation = float(el.get("rotate") or 0)
    setup_text_frame(box.text_frame, el, theme, emu_per_px)
    # PPTist 兼容：元素级 defaultColor/defaultFontName 作为 run 缺省
    el_defaults = {}
    if el.get("defaultColor"):
        el_defaults["color"] = el["defaultColor"]
    if el.get("defaultFontName"):
        el_defaults["fontName"] = el["defaultFontName"]
    render_paragraphs(box.text_frame, el.get("paragraphs") or [], theme, el.get("textType"), el_defaults)
    if el.get("fill"):
        apply_solid_fill(box.fill, el["fill"])
    else:
        box.fill.background()
    if el.get("outline"):
        apply_line_format(box.line, el["outline"], emu_per_px)
    else:
        box.line.fill.background()
    box.shadow.inherit = False
    if el.get("shadow"):
        apply_shadow(box._element.spPr, el["shadow"], emu_per_px)
    if el.get("opacity") is not None and float(el["opacity"]) < 1:
        # 文本框整体透明度：作用在 fill 上无意义时忽略（pptx 不支持文本框整体 alpha）
        pass
    return box


def resolve_image_src(src: str, deck_dir: Path) -> Path:
    if re.match(r"^https?://", src):
        _TMP_IMG_DIR.mkdir(parents=True, exist_ok=True)
        ext = os.path.splitext(src.split("?")[0])[1] or ".jpg"
        local = _TMP_IMG_DIR / (hashlib.md5(src.encode()).hexdigest() + ext)
        if not local.exists():
            req = urllib.request.Request(src, headers={"User-Agent": "Mozilla/5.0"})
            with urllib.request.urlopen(req, timeout=30) as resp, open(local, "wb") as f:
                f.write(resp.read())
        return local
    p = Path(src)
    if not p.is_absolute():
        p = deck_dir / p
    if not p.is_file():
        raise DeckError(f"图片不存在: {p}")
    return p


def _image_size(path: Path):
    if _PILImage is None:
        return None
    try:
        with _PILImage.open(path) as im:
            return im.size  # (w, h)
    except Exception:  # noqa: BLE001
        return None


def render_image(slide, el, deck_dir, emu_per_px):
    src = el.get("src") or ""
    if src.startswith("gen:") or src.startswith("icon:"):
        raise DeckError(
            f"图片占位符未解析: {src}（元素 {el.get('id')}）。请先运行 gen_image.py 填充图片。"
        )
    if not src:
        raise DeckError(f"image 元素 {el.get('id')} 缺少 src")
    path = resolve_image_src(src, deck_dir)
    left, top = px_to_emu(el["left"], emu_per_px), px_to_emu(el["top"], emu_per_px)
    w, h = px_to_emu(el["width"], emu_per_px), px_to_emu(el["height"], emu_per_px)
    fit = el.get("fit", "cover")
    size = _image_size(path)

    if fit == "contain" and size:
        iw, ih = size
        scale = min(w / iw, h / ih)
        dw, dh = int(iw * scale), int(ih * scale)
        pic = slide.shapes.add_picture(str(path), left + int((w - dw) / 2), top + int((h - dh) / 2), dw, dh)
    else:
        pic = slide.shapes.add_picture(str(path), left, top, w, h)
        if fit == "cover" and size:
            iw, ih = size
            target_ratio = w / h
            img_ratio = iw / ih
            if img_ratio > target_ratio:  # 图更宽 → 裁左右
                keep = target_ratio / img_ratio
                pic.crop_left = pic.crop_right = (1 - keep) / 2
            elif img_ratio < target_ratio:  # 图更高 → 裁上下
                keep = img_ratio / target_ratio
                pic.crop_top = pic.crop_bottom = (1 - keep) / 2
        # fit == "fill"：直接拉伸，无需裁剪

    if el.get("name"):
        pic.name = el["name"]
    pic.rotation = float(el.get("rotate") or 0)
    set_flip(pic, el.get("flipH"), el.get("flipV"))
    if el.get("radius"):
        set_picture_round(pic, float(el["radius"]), float(el["width"]), float(el["height"]))
    if el.get("outline"):
        apply_line_format(pic.line, el["outline"], emu_per_px)
    pic.shadow.inherit = False
    if el.get("shadow"):
        apply_shadow(pic._element.spPr, el["shadow"], emu_per_px)
    if el.get("opacity") is not None:
        set_picture_opacity(pic, float(el["opacity"]))
    return pic


def render_shape(slide, el, theme, emu_per_px):
    prst = el.get("shape", "rect")
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        px_to_emu(el["left"], emu_per_px), px_to_emu(el["top"], emu_per_px),
        px_to_emu(el["width"], emu_per_px), px_to_emu(el["height"], emu_per_px),
    )
    set_shape_geometry(shape, prst, el.get("adjust"))
    # 形状名优先取 id：PowerPoint 选择窗格里能按元素 id 找到它，方便二次编辑
    shape.name = el.get("name") or el.get("id") or shape.name
    shape.rotation = float(el.get("rotate") or 0)
    set_flip(shape, el.get("flipH"), el.get("flipV"))

    if el.get("gradient"):
        apply_gradient_fill(shape.fill, el["gradient"], emu_per_px)
    elif el.get("fill"):
        rgb, alpha = parse_color(el["fill"])
        if el.get("opacity") is not None and float(el["opacity"]) < 1:
            alpha = float(el["opacity"]) if alpha is None else alpha * float(el["opacity"])
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb
        if alpha is not None:
            solid = shape.fill._xPr.find(qn("a:solidFill"))
            if solid is not None:
                _set_solid_fill_alpha(solid, alpha)
    else:
        shape.fill.background()

    apply_line_format(shape.line, el.get("outline"), emu_per_px)
    shape.shadow.inherit = False
    if el.get("shadow"):
        apply_shadow(shape._element.spPr, el["shadow"], emu_per_px)

    text_spec = el.get("text")
    if text_spec and text_spec.get("paragraphs"):
        setup_text_frame(shape.text_frame, text_spec, theme, emu_per_px)
        el_defaults = {}
        if text_spec.get("defaultColor"):
            el_defaults["color"] = text_spec["defaultColor"]
        if text_spec.get("defaultFontName"):
            el_defaults["fontName"] = text_spec["defaultFontName"]
        render_paragraphs(shape.text_frame, text_spec["paragraphs"], theme, text_spec.get("textType"), el_defaults)
    else:
        shape.text_frame.word_wrap = True
    return shape


def render_line(slide, el, emu_per_px):
    sx, sy = el["start"]
    ex, ey = el["end"]
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        px_to_emu(sx, emu_per_px), px_to_emu(sy, emu_per_px),
        px_to_emu(ex, emu_per_px), px_to_emu(ey, emu_per_px),
    )
    if el.get("name"):
        conn.name = el["name"]
    apply_line_format(
        conn.line,
        {"style": el.get("style", "solid"), "width": el.get("width", 2), "color": el.get("color", "#1F2937")},
        emu_per_px,
    )
    set_arrowheads(conn, el.get("points") or ["", ""])
    conn.shadow.inherit = False
    if el.get("shadow"):
        apply_shadow(conn._element.spPr, el["shadow"], emu_per_px)
    return conn


def render_chart(slide, el, theme, emu_per_px):
    ctype = el.get("chartType", "column")
    stack = bool((el.get("options") or {}).get("stack"))
    smooth = bool((el.get("options") or {}).get("lineSmooth"))
    data = el["data"]
    labels = [str(x) for x in data.get("labels") or []]
    legends = [str(x) for x in data.get("legends") or []]
    series = data.get("series") or []
    colors = el.get("themeColors") or theme.get("themeColors") or DEFAULT_THEME["themeColors"]

    xl_map = {
        "column": XL_CHART_TYPE.COLUMN_STACKED if stack else XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar": XL_CHART_TYPE.BAR_STACKED if stack else XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS_STACKED if stack else XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "ring": XL_CHART_TYPE.DOUGHNUT,
        "area": XL_CHART_TYPE.AREA_STACKED if stack else XL_CHART_TYPE.AREA,
        "radar": XL_CHART_TYPE.RADAR_MARKERS,
    }
    frame = None
    if ctype == "scatter":
        chart_data = XyChartData()
        for i, ys in enumerate(series):
            s = chart_data.add_series(legends[i] if i < len(legends) else f"系列{i + 1}")
            for x, y in zip(labels, ys):
                s.add_data_point(float(x), float(y))
        frame = slide.shapes.add_chart(
            XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS if smooth else XL_CHART_TYPE.XY_SCATTER_LINES,
            px_to_emu(el["left"], emu_per_px), px_to_emu(el["top"], emu_per_px),
            px_to_emu(el["width"], emu_per_px), px_to_emu(el["height"], emu_per_px),
            chart_data,
        )
    else:
        chart_data = CategoryChartData()
        chart_data.categories = labels
        for i, ys in enumerate(series):
            chart_data.add_series(legends[i] if i < len(legends) else f"系列{i + 1}", [float(v) for v in ys])
        frame = slide.shapes.add_chart(
            xl_map[ctype],
            px_to_emu(el["left"], emu_per_px), px_to_emu(el["top"], emu_per_px),
            px_to_emu(el["width"], emu_per_px), px_to_emu(el["height"], emu_per_px),
            chart_data,
        )
    chart = frame.chart
    chart.has_title = False
    if ctype not in ("pie", "ring"):
        chart.has_legend = len(series) > 1
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
    elif len(legends) > 1 and ctype in ("pie", "ring"):
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False

    # 主题色
    for i, plot_series in enumerate(chart.series):
        rgb, _ = parse_color(colors[i % len(colors)])
        plot_series.format.fill.solid()
        plot_series.format.fill.fore_color.rgb = rgb
        if ctype in ("line", "scatter", "radar"):
            plot_series.format.line.color.rgb = rgb
            plot_series.format.line.width = Pt(2)
        if smooth and ctype == "line":
            plot_series.smooth = True
        if ctype in ("pie", "ring"):
            # 饼/环按数据点取色
            for j, point in enumerate(plot_series.points):
                prgb, _ = parse_color(colors[j % len(colors)])
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = prgb
    if el.get("name"):
        frame.name = el["name"]
    return frame


def _set_cell_border(cell, outline: dict, emu_per_px: float):
    rgb, _ = parse_color(outline.get("color", "#E5E7EB"))
    w = str(int(outline.get("width", 1) * emu_per_px))
    dash = DASH_MAP.get(outline.get("style", "solid"))
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for old in tcPr.findall(qn(tag)):
            tcPr.remove(old)
        ln = _sub(tcPr, tag, w=w, cap="flat", cmpd="sng", algn="ctr")
        solid = _sub(ln, "a:solidFill")
        _sub(solid, "a:srgbClr", val=str(rgb))
        if dash:
            _sub(ln, "a:prstDash", val=dash)


def render_table(slide, el, theme, emu_per_px):
    data = el.get("data") or []
    if not data:
        raise DeckError(f"table 元素 {el.get('id')} data 为空")
    rows, cols = len(data), max(len(r) for r in data)
    col_widths = el.get("colWidths") or [1.0 / cols] * cols
    frame = slide.shapes.add_table(
        rows, cols,
        px_to_emu(el["left"], emu_per_px), px_to_emu(el["top"], emu_per_px),
        px_to_emu(el["width"], emu_per_px), px_to_emu(el["height"], emu_per_px),
    )
    if el.get("name"):
        frame.name = el["name"]
    table = frame.table
    table.first_row = False
    table.horz_banding = False

    total_w = float(el["width"])
    for i in range(cols):
        frac = col_widths[i] if i < len(col_widths) else 1.0 / cols
        table.columns[i].width = px_to_emu(total_w * frac, emu_per_px)
    min_h = float(el.get("cellMinHeight", 36))
    for r in range(rows):
        table.rows[r].height = px_to_emu(min_h, emu_per_px)

    tbl_theme = el.get("theme") or {}
    header_color = tbl_theme.get("color") or (theme.get("themeColors") or ["#1E40AF"])[0]
    header_rgb, _ = parse_color(header_color)
    header_text_color = contrast_text_color(header_rgb)
    default_outline = (theme.get("outline") or {}).get("color", "#E5E7EB")
    outline = el.get("outline") or {"style": "solid", "width": 1, "color": default_outline}
    # 表体填充跟随主题背景：深底提亮一档作卡片面，浅底用纯白；避免深主题下白底+浅字不可读
    bg_rgb, _ = parse_color(theme.get("backgroundColor", "#FFFFFF"))
    body_fill = shift_color(bg_rgb, 0.10) if rel_luminance(bg_rgb) < 0.35 else RGBColor(0xFF, 0xFF, 0xFF)

    for r, row in enumerate(data):
        for c in range(cols):
            spec = row[c] if c < len(row) else {"text": "", "merged": True}
            cell = table.cell(r, c)
            if spec.get("merged"):
                continue
            style = spec.get("style") or {}
            is_header = bool(tbl_theme.get("rowHeader")) and r == 0
            is_footer = bool(tbl_theme.get("rowFooter")) and r == rows - 1
            is_col_header = bool(tbl_theme.get("colHeader")) and c == 0

            # 填充
            backcolor = style.get("backcolor")
            if backcolor:
                fill_color, alpha = parse_color(backcolor)
            elif is_header or is_footer or is_col_header:
                fill_color, alpha = parse_color(header_color)
            else:
                fill_color, alpha = body_fill, None
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color

            # 文本
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_top = tf.margin_bottom = px_to_emu(4, emu_per_px)
            tf.margin_left = tf.margin_right = px_to_emu(8, emu_per_px)
            cell.vertical_anchor = VALIGN_MAP.get(style.get("vAlign", "middle"), MSO_ANCHOR.MIDDLE)
            p = tf.paragraphs[0]
            p.alignment = ALIGN_MAP.get(style.get("align", "left"), PP_ALIGN.LEFT)
            run = p.add_run()
            run.text = str(spec.get("text", ""))
            f = run.font
            f.size = px_to_pt(float(_pick(style, "fontSize", "fontsize") or 13))
            emphasized = is_header or is_footer or is_col_header
            f.bold = bool(style.get("bold", emphasized))
            f.italic = bool(style.get("em", False))
            if style.get("underline"):
                f.underline = True
            default_color = header_text_color if (emphasized and not style.get("backcolor")) else theme.get("fontColor", "#1F2937")
            rgb, _ = parse_color(style.get("color") or default_color)
            f.color.rgb = rgb
            cell_font = _pick(style, "fontName", "fontname")
            set_run_font(run, cell_font or theme.get("fontName", "Microsoft YaHei"),
                         None if cell_font else theme.get("latinFontName"))
            if style.get("strikethrough"):
                set_strikethrough(run)
            _set_cell_border(cell, outline, emu_per_px)

            # 合并（在样式设置之后执行）
            colspan = int(spec.get("colspan") or 1)
            rowspan = int(spec.get("rowspan") or 1)
            if colspan > 1 or rowspan > 1:
                end = table.cell(min(r + rowspan - 1, rows - 1), min(c + colspan - 1, cols - 1))
                cell.merge(end)
    return frame


# ---------------------------------------------------------------------------
# LaTeX → OMML（PowerPoint 原生公式，可在 PPT 里继续编辑）
# 链路: latex2mathml (LaTeX→MathML) → mathml2omml (MathML→OMML) → 注入 a:p
# ---------------------------------------------------------------------------

def latex_to_omml(latex_str: str) -> str:
    from latex2mathml.converter import convert as _l2m
    from mathml2omml import convert as _m2o

    mathml = _l2m(latex_str)
    return _m2o(mathml)


def _unwrap_math_boxes(omath):
    """mathml2omml 用 m:box 做分组容器；m:box 默认带边框，拆箱还原为纯分组。"""
    for box in list(omath.iter(qn("m:box"))):
        parent = box.getparent()
        if parent is None:
            continue
        idx = list(parent).index(box)
        e = box.find(qn("m:e"))
        children = list(e) if e is not None else []
        parent.remove(box)
        for i, child in enumerate(children):
            parent.insert(idx + i, child)


# Cambria Math 字体声明（Mac Office 原生公式的固定写法，缺了 Mac 会静默丢弃 math zone）
_MATH_LATIN = {"typeface": "Cambria Math", "panose": "02040503050406030204",
               "pitchFamily": "18", "charset": "0"}


def _math_rpr(font_size_px: float, color_value, italic: bool):
    """构造 Mac Office 原生风格的 a:rPr：sz + i + solidFill + Cambria Math。"""
    from lxml import etree

    rgb, alpha = parse_color(color_value)
    sz = str(int(round(font_size_px * PT_PER_PX * 100)))
    attrs = {"sz": sz}
    if italic:
        attrs["i"] = "1"
    rPr = etree.Element(qn("a:rPr"), attrs)
    solid = _sub(rPr, "a:solidFill")
    clr = _sub(solid, "a:srgbClr", val=str(rgb))
    if alpha is not None:
        _sub(clr, "a:alpha", val=str(int(alpha * 100000)))
    _sub(rPr, "a:latin", **_MATH_LATIN)
    return rPr


# 结构元素 → 其 *Pr 子元素（Mac 原生在每个 *Pr 里都带 m:ctrlPr）
_MATH_STRUCT_PR = {
    "m:rad": "m:radPr", "m:f": "m:fPr", "m:sSub": "m:sSubPr", "m:sSup": "m:sSupPr",
    "m:sSubSup": "m:sSubSupPr", "m:d": "m:dPr", "m:nary": "m:naryPr",
    "m:func": "m:funcPr", "m:groupChr": "m:groupChrPr", "m:limLow": "m:limLowPr",
    "m:limUpp": "m:limUppPr", "m:bar": "m:barPr", "m:m": "m:mPr", "m:eqArr": "m:eqArrPr",
}
_MATH_STRUCT_PR_QN = {qn(k): qn(v) for k, v in _MATH_STRUCT_PR.items()}
_QN_MR, _QN_MT, _QN_MRPR, _QN_MSTY = qn("m:r"), qn("m:t"), qn("m:rPr"), qn("m:sty")
_QN_CTRLPR, _QN_RAD, _QN_DEGHIDE, _QN_DEG, _QN_E = (
    qn("m:ctrlPr"), qn("m:rad"), qn("m:degHide"), qn("m:deg"), qn("m:e"))
_QN_VAL = qn("m:val")


def _macify_omml(omath, font_size_px: float, color_value):
    """把 mathml2omml 的输出改写成 Mac PowerPoint 原生公式的字节模式。

    实测对照（Mac Office 365 插入公式后存盘）：原生公式没有 m:rPr/m:sty，
    样式全在 a:rPr（i="1" 表斜体），且每个 run 都显式声明 Cambria Math；
    每个结构元素（rad/f/sSub…）的 *Pr 里都有 m:ctrlPr，m:rad 还带空 m:deg。
    mathml2omml 的输出缺这些，Mac 解析器拿到无法排版的 math zone 会静默丢成空盒子。

    单次深度优先遍历完成全部改写（原实现对 15 类标签各扫一遍全树，O(15N)）。
    """
    for el in list(omath.iter()):  # list() 快照：遍历时增删子元素不影响迭代
        tag = el.tag
        if tag == _QN_MR:
            t = el.find(_QN_MT)
            if t is None:
                continue
            mrpr = el.find(_QN_MRPR)
            italic = False
            if mrpr is not None:
                sty = mrpr.find(_QN_MSTY)
                italic = sty is not None and sty.get(_QN_VAL) == "i"
                el.remove(mrpr)  # Mac 原生没有 m:rPr，斜体走 a:rPr 的 i 属性
            el.insert(list(el).index(t), _math_rpr(font_size_px, color_value, italic))
            continue
        pr_tag = _MATH_STRUCT_PR_QN.get(tag)
        if pr_tag is None:
            continue
        pr = el.find(pr_tag)
        if pr is None:
            pr = el.makeelement(pr_tag, {})
            el.insert(0, pr)
        if pr.find(_QN_CTRLPR) is None:
            ctrl = pr.makeelement(_QN_CTRLPR, {})
            ctrl.append(_math_rpr(font_size_px, color_value, True))
            pr.append(ctrl)
        # m:rad 的子元素顺序是 radPr? deg? e —— Mac 原生带空 m:deg + degHide on，
        # 缺了 Mac 不渲染根号；degHide 隐藏空次数占位框（\sqrt 无次数）
        if tag == _QN_RAD:
            if pr.find(_QN_DEGHIDE) is None:
                pr.append(pr.makeelement(_QN_DEGHIDE, {_QN_VAL: "on"}))  # ctrlPr 前、degHide 后
            if el.find(_QN_DEG) is None:
                e = el.find(_QN_E)
                deg = el.makeelement(_QN_DEG, {})
                el.insert(list(el).index(e) if e is not None else len(el), deg)


# ---------------------------------------------------------------------------
# LaTeX → 原生文本 run（备选引擎）：把公式转成 PPT 普通文本 + 真实上下标 run。
# 默认引擎是 OMML（见上，照 Mac Office 原生格式注入，真根号/真分式、可编辑）；
# 但 Keynote 不支持 a14:m（会判整个文件非法），需要 Keynote 交付时在 latex
# 元素上加 "engine": "runs"——代价是根号没有横线、分式是行内式 (a)/(b)。
# ---------------------------------------------------------------------------

_LATEX_GREEK = {
    "alpha": "α", "beta": "β", "gamma": "γ", "delta": "δ", "epsilon": "ε",
    "zeta": "ζ", "eta": "η", "theta": "θ", "iota": "ι", "kappa": "κ",
    "lambda": "λ", "mu": "μ", "nu": "ν", "xi": "ξ", "pi": "π", "rho": "ρ",
    "sigma": "σ", "tau": "τ", "phi": "φ", "chi": "χ", "psi": "ψ", "omega": "ω",
    "Gamma": "Γ", "Delta": "Δ", "Theta": "Θ", "Lambda": "Λ", "Pi": "Π",
    "Sigma": "Σ", "Phi": "Φ", "Omega": "Ω",
}
_LATEX_SYMBOLS = {
    "cdot": "·", "times": "×", "otimes": "⊗", "oplus": "⊕", "infty": "∞",
    "pm": "±", "mp": "∓", "leq": "≤", "geq": "≥", "neq": "≠", "approx": "≈",
    "partial": "∂", "sum": "Σ", "prod": "∏", "int": "∫", "nabla": "∇",
    "rightarrow": "→", "to": "→", "leftarrow": "←", "Rightarrow": "⇒",
    "in": "∈", "subset": "⊂", "cup": "∪", "cap": "∩", "forall": "∀", "exists": "∃",
    "quad": "  ", "qquad": "    ", ",": " ", ";": " ", " ": " ", "!": "",
}
_LATEX_WORDS = {
    "sin", "cos", "tan", "log", "ln", "exp", "max", "min", "argmax", "argmin",
    "softmax", "sigmoid", "tanh", "det", "dim", "ker", "deg", "gcd",
}


def _latex_tokenize(src: str):
    toks, i = [], 0
    while i < len(src):
        c = src[i]
        if c == "\\":
            j = i + 1
            if j < len(src) and src[j].isalpha():
                k = j
                while k < len(src) and src[k].isalpha():
                    k += 1
                toks.append(("cmd", src[j:k]))
                i = k
            else:  # 单字符命令（\, \; \{ 等）
                toks.append(("cmd", src[j:j + 1]))
                i = j + 1
        elif c == "{":
            toks.append(("lb", c)); i += 1
        elif c == "}":
            toks.append(("rb", c)); i += 1
        elif c == "^":
            toks.append(("sup", c)); i += 1
        elif c == "_":
            toks.append(("sub", c)); i += 1
        elif c == " ":
            toks.append(("sp", c)); i += 1
        else:
            toks.append(("ch", c)); i += 1
    return toks


def _latex_parse_arg(toks, i):
    """解析 ^/_ 或 \sqrt/\frac 的参数：{组} | 单个命令 | 单个字符。"""
    if i < len(toks) and toks[i][0] == "lb":
        return _latex_parse_seq(toks, i + 1)
    if i < len(toks) and toks[i][0] == "cmd":
        runs, ni = _latex_cmd_runs(toks, i)
        return runs, ni
    if i < len(toks) and toks[i][0] == "ch":
        return [{"text": toks[i][1]}], i + 1
    return [], i


def _latex_cmd_runs(toks, i):
    """把一个命令（及其参数）转成 run 列表。返回 (runs, next_i)。"""
    name = toks[i][1]
    i += 1
    if name in ("mathrm", "text", "mathbf", "operatorname", "mathit", "mathcal"):
        return _latex_parse_arg(toks, i)
    if name == "frac" or name == "dfrac" or name == "tfrac":
        num, i = _latex_parse_arg(toks, i)
        den, i = _latex_parse_arg(toks, i)
        num_t = "".join(r["text"] for r in num)
        den_t = "".join(r["text"] for r in den)
        out = []
        if len(num_t) > 1:
            out.append({"text": "("})
        out.extend(num)
        if len(num_t) > 1:
            out.append({"text": ")"})
        out.append({"text": "/"})
        if len(den_t) > 1 or not den:
            out.append({"text": "("})
        out.extend(den)
        if len(den_t) > 1 or not den:
            out.append({"text": ")"})
        return out, i
    if name == "sqrt":
        arg, i = _latex_parse_arg(toks, i)
        arg_t = "".join(r["text"] for r in arg)
        out = [{"text": "√"}]
        if any(c in arg_t for c in "+-/=·× ") or len(arg_t) > 3:
            out.append({"text": "("})
            out.extend(arg)
            out.append({"text": ")"})
        else:
            out.extend(arg)
        return out, i
    if name in ("left", "right", "limits", "displaystyle", "big", "Big"):
        return [], i
    if name in _LATEX_GREEK:
        return [{"text": _LATEX_GREEK[name]}], i
    if name in _LATEX_SYMBOLS:
        return [{"text": _LATEX_SYMBOLS[name]}], i
    if name in _LATEX_WORDS:
        return [{"text": name}], i
    if name == "{":
        return [{"text": "{"}], i
    if name == "}":
        return [{"text": "}"}], i
    return [{"text": name}], i  # 未知命令：保留名字保证可读


def _latex_parse_seq(toks, i):
    runs = []
    while i < len(toks):
        kind, val = toks[i]
        if kind == "rb":
            return runs, i + 1
        if kind == "lb":
            sub, i = _latex_parse_seq(toks, i + 1)
            runs.extend(sub)
            continue
        if kind in ("sup", "sub"):
            arg, i = _latex_parse_arg(toks, i + 1)
            for r in arg:
                r[kind] = True
            runs.extend(arg)
            continue
        if kind == "cmd":
            out, i = _latex_cmd_runs(toks, i)
            runs.extend(out)
            continue
        runs.append({"text": " " if kind == "sp" else val})
        i += 1
    return runs, i


def _latex_to_runs(latex_src: str):
    """LaTeX → run 列表（相邻同样式合并），每个 run: {text, sub?, sup?}。"""
    runs, _ = _latex_parse_seq(_latex_tokenize(latex_src), 0)
    merged = []
    for r in runs:
        if merged and merged[-1].get("sub") == r.get("sub") and merged[-1].get("sup") == r.get("sup"):
            merged[-1]["text"] += r["text"]
        else:
            merged.append(dict(r))
    return merged


def render_latex(slide, el, theme, emu_per_px):
    if (el.get("engine") or "omml") == "omml":
        return _render_latex_omml(slide, el, theme, emu_per_px)
    latex_src = str(el.get("latex", ""))
    font_size = float(el.get("fontSize") or 20)
    color = el.get("color") or theme.get("fontColor", "#1F2937")
    # 两分支共用的基础字段构建一次，只在 paragraphs 上分叉（否则降级分支会静默丢 rotate 等字段）
    base_text_el = {
        "id": el.get("id"), "name": el.get("name"),
        "left": el["left"], "top": el["top"],
        "width": el["width"], "height": el["height"],
        "rotate": el.get("rotate"),
        "text": el.get("text"),
    }
    try:
        runs = _latex_to_runs(latex_src)
        if not runs:
            raise ValueError("空公式")
        for r in runs:
            r["fontSize"] = font_size
            r["color"] = color
        text_el = {
            **base_text_el,
            "paragraphs": [{"align": el.get("align", "center"), "runs": runs, "lineHeight": 1.3}],
        }
        box = render_text(slide, text_el, theme, emu_per_px)
        # LaTeX 源码存进选择窗格描述，后续可对照手动编辑
        box._element.nvSpPr.cNvPr.set("descr", latex_src)
        return box
    except Exception as e:  # noqa: BLE001 - 降级为原文文本，不阻断整页渲染
        print(f"[warn] latex 元素 {el.get('id')} 转换失败（{e}），已降级为源码文本", file=sys.stderr)
        text_el = {
            **base_text_el,
            "paragraphs": [{"align": el.get("align", "center"),
                            "runs": [{"text": latex_src, "fontSize": font_size, "color": color}]}],
        }
        return render_text(slide, text_el, theme, emu_per_px)


def _render_latex_omml(slide, el, theme, emu_per_px):
    box = slide.shapes.add_textbox(
        px_to_emu(el["left"], emu_per_px), px_to_emu(el["top"], emu_per_px),
        px_to_emu(el["width"], emu_per_px), px_to_emu(el["height"], emu_per_px),
    )
    if el.get("name") or el.get("id"):
        box.name = el.get("name") or el["id"]
    box.rotation = float(el.get("rotate") or 0)
    setup_text_frame(box.text_frame, el, theme, emu_per_px)
    box.fill.background()
    box.line.fill.background()
    box.shadow.inherit = False

    p = box.text_frame.paragraphs[0]
    p.alignment = ALIGN_MAP.get(el.get("align", "center"), PP_ALIGN.CENTER)
    latex_src = str(el.get("latex", ""))
    font_size = float(el.get("fontSize") or 20)
    color = el.get("color") or theme.get("fontColor", "#1F2937")
    try:
        from lxml import etree

        etree.register_namespace("m", M_NS)
        etree.register_namespace("a14", A14_NS)
        omml_str = latex_to_omml(latex_src)
        wrapped = f'<root xmlns:m="{M_NS}" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">{omml_str}</root>'
        root = etree.fromstring(wrapped.encode("utf-8"))
        omath = root[0]
        _unwrap_math_boxes(omath)
        _macify_omml(omath, font_size, color)
        # 对照 Mac Office 365 真实存盘：公式就是裸 a14:m > m:oMathPara，没有
        # mc:AlternateContent——Mac 遇到 Requires="a14" 的 Choice 会直接选
        # Fallback 显示纯文本，所以这里必须照原生格式裸注入。
        jc = {"left": "left", "center": "center", "right": "right"}.get(el.get("align", "center"), "center")
        omath_para = omath.makeelement(f"{{{M_NS}}}oMathPara", {})
        para_pr = omath.makeelement(f"{{{M_NS}}}oMathParaPr", {})
        jc_el = omath.makeelement(f"{{{M_NS}}}jc", {f"{{{M_NS}}}val": jc})
        para_pr.append(jc_el)
        omath_para.append(para_pr)
        omath_para.append(omath)
        a14_m = p._p.makeelement(f"{{{A14_NS}}}m", {})
        a14_m.append(omath_para)
        p._p.append(a14_m)
    except Exception as e:  # noqa: BLE001 - 降级为原文文本，不阻断整页渲染
        print(f"[warn] latex 元素 {el.get('id')} 转换失败（{e}），已降级为纯文本", file=sys.stderr)
        run = p.add_run()
        run.text = latex_src
        run.font.size = px_to_pt(font_size)
        run.font.color.rgb = parse_color(color)[0]
        set_run_font(run, theme.get("fontName", "Microsoft YaHei"))
    return box


def render_background(slide, bg: dict | None, theme, canvas_w, canvas_h, emu_per_px, deck_dir):
    bg = bg or {"type": "solid", "color": theme.get("backgroundColor", "#FFFFFF")}
    btype = bg.get("type", "solid")
    if btype == "gradient" and bg.get("gradient"):
        apply_gradient_fill(slide.background.fill, bg["gradient"], emu_per_px)
    elif btype == "image" and (bg.get("image") or {}).get("src"):
        # 底色先铺上，图片以全幅 cover 插入到最底层（第一个添加）
        color = bg.get("color") or theme.get("backgroundColor", "#FFFFFF")
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = parse_color(color)[0]
        img = bg["image"]
        el = {
            "id": "__bg_image__", "type": "image", "left": 0, "top": 0,
            "width": canvas_w, "height": canvas_h,
            "src": img["src"], "fit": "cover" if img.get("size", "cover") == "cover" else "contain",
        }
        render_image(slide, el, deck_dir, emu_per_px)
    else:
        color = bg.get("color") or theme.get("backgroundColor", "#FFFFFF")
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = parse_color(color)[0]


# ---------------------------------------------------------------------------
# 校验
# ---------------------------------------------------------------------------

def validate_deck(deck: dict):
    errors, warnings = [], []
    if not isinstance(deck.get("slides"), list) or not deck["slides"]:
        errors.append("deck.slides 必须是非空数组")
        return errors, warnings
    canvas = deck.get("canvas") or {}
    cw, ch = float(canvas.get("width", DEFAULT_CANVAS_W)), float(canvas.get("height", DEFAULT_CANVAS_H))
    seen_slide_ids = set()

    for si, slide in enumerate(deck["slides"]):
        sprefix = f"slides[{si}]"
        sid = slide.get("id")
        if not sid:
            errors.append(f"{sprefix} 缺少 id")
        elif sid in seen_slide_ids:
            errors.append(f"{sprefix} id 重复: {sid}")
        seen_slide_ids.add(sid)
        seen_el_ids = set()  # 元素 id 只需页内唯一
        stype = slide.get("type")
        if stype and stype not in SLIDE_TYPES:
            warnings.append(f"{sprefix} 未知 slideType: {stype}")
        elements = slide.get("elements")
        if not isinstance(elements, list):
            errors.append(f"{sprefix} 缺少 elements 数组")
            continue
        for ei, el in enumerate(elements):
            ep = f"{sprefix}.elements[{ei}]({el.get('id', '?')})"
            etype = el.get("type")
            if etype not in ("text", "image", "shape", "line", "chart", "table", "latex"):
                errors.append(f"{ep} 未知元素类型: {etype}")
                continue
            eid = el.get("id")
            if not eid:
                errors.append(f"{ep} 缺少 id")
            elif eid in seen_el_ids:
                warnings.append(f"{ep} 元素 id 重复: {eid}")
            seen_el_ids.add(eid)

            if etype == "line":
                if not el.get("start") or not el.get("end"):
                    errors.append(f"{ep} line 需要 start/end")
            else:
                for k in ("left", "top", "width", "height"):
                    if not isinstance(el.get(k), (int, float)):
                        errors.append(f"{ep} 缺少数值字段 {k}")
                else:
                    if el.get("left", 0) < -1 or el.get("top", 0) < -1:
                        warnings.append(f"{ep} 位置为负（超出画布左上）")
                    if isinstance(el.get("left"), (int, float)) and isinstance(el.get("width"), (int, float)):
                        if el["left"] + el["width"] > cw + 1 or el.get("top", 0) + el.get("height", 0) > ch + 1:
                            warnings.append(f"{ep} 超出画布右/下边界")

            if etype == "text":
                if not el.get("paragraphs"):
                    errors.append(f"{ep} text 缺少 paragraphs")
                tt = el.get("textType")
                if tt and tt not in TEXT_TYPES:
                    warnings.append(f"{ep} 未知 textType: {tt}")
            elif etype == "latex":
                if not el.get("latex"):
                    errors.append(f"{ep} latex 元素缺少 latex 字段（公式源码）")
            elif etype == "image":
                src = el.get("src") or ""
                if src.startswith("gen:") or src.startswith("icon:"):
                    errors.append(f"{ep} 图片占位符未解析: {src}，请先运行 gen_image.py")
                elif not src:
                    warnings.append(f"{ep} image 缺少 src（如需自动生成请用 gen:关键词 或 icon:集合:名称）")
                it = el.get("imageType")
                if it and it not in IMAGE_TYPES:
                    warnings.append(f"{ep} 未知 imageType: {it}")
            elif etype == "shape":
                prst = el.get("shape", "rect")
                if prst not in PRESET_SHAPES:
                    errors.append(f"{ep} 未知预设形状: {prst}（可用: {', '.join(sorted(PRESET_SHAPES))}）")
            elif etype == "chart":
                ct = el.get("chartType")
                if ct not in CHART_TYPES:
                    errors.append(f"{ep} 未知 chartType: {ct}")
                data = el.get("data") or {}
                labels = data.get("labels") or []
                for i, ys in enumerate(data.get("series") or []):
                    if len(ys) != len(labels):
                        errors.append(f"{ep} series[{i}] 长度 {len(ys)} 与 labels 长度 {len(labels)} 不一致")
            elif etype == "table":
                rows = el.get("data") or []
                if not rows:
                    errors.append(f"{ep} table data 为空")
                else:
                    ncol = max(len(r) for r in rows)
                    for ri, r in enumerate(rows):
                        if len(r) != ncol:
                            warnings.append(f"{ep} 第 {ri} 行列数 {len(r)} 与最多列数 {ncol} 不一致")
                    cw_sum = sum(el.get("colWidths") or [])
                    if el.get("colWidths") and abs(cw_sum - 1.0) > 0.02:
                        warnings.append(f"{ep} colWidths 之和为 {cw_sum:.3f}（建议归一到 1.0）")
    return errors, warnings


# ---------------------------------------------------------------------------
# 主流程
# ---------------------------------------------------------------------------

def render_deck(deck: dict, out_path: Path, theme: dict, deck_dir: Path):
    canvas = deck.get("canvas") or {}
    cw = float(canvas.get("width", DEFAULT_CANVAS_W))
    ch = float(canvas.get("height", DEFAULT_CANVAS_H))
    emu_per_px = SLIDE_W_EMU / cw
    slide_h_emu = int(round(ch * emu_per_px))

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(slide_h_emu)
    blank = prs.slide_layouts[6]

    stats = {"slides": 0, "elements": 0}
    for slide_spec in deck["slides"]:
        slide = prs.slides.add_slide(blank)
        render_background(slide, slide_spec.get("background"), theme, cw, ch, emu_per_px, deck_dir)
        for el in slide_spec.get("elements") or []:
            etype = el["type"]
            if etype == "text":
                render_text(slide, el, theme, emu_per_px)
            elif etype == "image":
                render_image(slide, el, deck_dir, emu_per_px)
            elif etype == "shape":
                render_shape(slide, el, theme, emu_per_px)
            elif etype == "line":
                render_line(slide, el, emu_per_px)
            elif etype == "chart":
                render_chart(slide, el, theme, emu_per_px)
            elif etype == "table":
                render_table(slide, el, theme, emu_per_px)
            elif etype == "latex":
                render_latex(slide, el, theme, emu_per_px)
            stats["elements"] += 1
        if slide_spec.get("remark"):
            slide.notes_slide.notes_text_frame.text = str(slide_spec["remark"])
        stats["slides"] += 1

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return stats


def main():
    ap = argparse.ArgumentParser(description="瘦身版 PPTist JSON → pptx 渲染器")
    ap.add_argument("deck", help="deck JSON 文件路径，或项目目录（含 deck.json + pages/*.json）")
    ap.add_argument("-o", "--out", help="输出 pptx 路径（默认 <deck>.pptx；项目目录时为 <目录名>.pptx）")
    ap.add_argument("--theme", help="覆盖 deck 里的主题（scripts/themes/ 下的主题名）")
    ap.add_argument("--check", action="store_true", help="只校验不渲染")
    args = ap.parse_args()

    deck_path = Path(args.deck).resolve()
    deck, deck_dir, _page_files = load_deck(deck_path)
    script_dir = Path(__file__).resolve().parent

    errors, warnings = validate_deck(deck)
    for w in warnings:
        print(f"[warn] {w}", file=sys.stderr)
    if errors:
        for e in errors:
            print(f"[error] {e}", file=sys.stderr)
        sys.exit(1)
    if args.check:
        print(f"[ok] 校验通过（{len(deck['slides'])} 页，{len(warnings)} 个警告）")
        return

    theme = load_theme(deck, args.theme, script_dir)
    if args.out:
        out = Path(args.out).resolve()
    elif deck_path.is_dir():
        out = deck_path / (deck_path.name + ".pptx")
    else:
        out = deck_path.with_suffix(".pptx")
    stats = render_deck(deck, out, theme, deck_dir)
    print(f"[ok] 已生成 {out}（{stats['slides']} 页 / {stats['elements']} 个元素，主题 {theme.get('name', 'custom')}）")


if __name__ == "__main__":
    main()
