"""render_pptx 文字溢出估算与 autofit 兜底的回归测试。

覆盖：字宽系数（中西文分开）、换行模拟（拉丁词不拆/CJK 逐字/行尾空格不计）、
行高模型（单倍 ≈1.32×字号，显式 lineHeight 再乘）、三级溢出判定（紧凑单行条
放行 / fixable warn / unfixable error / autoFit=none error）、表格行高估算与
两级表格检查（超设定高 warn、超画布 error）、normAutofit XML 写入（显式
fontScale + 剥掉 add_textbox 模板自带的 spAutoFit）。

依赖 python-pptx（render_pptx 运行时依赖，未进项目 venv，缺则 skip）。
"""
import importlib.util
import json
import subprocess
import sys
from pathlib import Path

import pytest

pytest.importorskip("pptx")

SCRIPTS = Path(__file__).resolve().parent.parent / "ethan/defaults/skills/ppt-generate/scripts"
sys.path.insert(0, str(SCRIPTS))

spec = importlib.util.spec_from_file_location("render_pptx", SCRIPTS / "render_pptx.py")
render_pptx = importlib.util.module_from_spec(spec)
sys.modules["render_pptx"] = render_pptx
spec.loader.exec_module(render_pptx)

THEME = render_pptx.DEFAULT_THEME
QN = render_pptx.qn


def _content_defaults():
    return render_pptx.typo_default(THEME, "content")


def _issues_for(*elements):
    deck = {"slides": [{"id": "s1", "elements": list(elements)}]}
    return render_pptx.validate_deck(deck, THEME)


def _codes(issues):
    return [i["code"] for i in issues if i["code"]]


# --- 字宽与换行 -------------------------------------------------------------


def test_char_widths():
    f = render_pptx.char_w_em
    assert f("中") == 1.0 and f("，") == 1.0 and f("。") == 1.0
    assert f("“") == 1.0 and f("…") == 1.0  # 全宽标点（不在 U+3000 段内）
    assert f("A") == 0.72 and f("A", bold=True) == 0.75
    assert f("a") == 0.55 and f("a", bold=True) == 0.58
    assert f("5") == 0.60 and f(" ") == 0.32


def test_wrap_cjk_by_char():
    """30 个中文字 @14px 在 250px 宽 → 17+13 两行。"""
    atoms = render_pptx._tokenize_runs([{"text": "字" * 30, "fontSize": 14}], _content_defaults())
    lines = render_pptx._wrap_lines(atoms, 250, 250)
    assert len(lines) == 2
    assert lines[0] == 14.0  # 每行最大字号


def test_wrap_latin_word_not_split():
    """拉丁词是原子，放不下也不拆（对齐 PowerPoint 断词换行）。"""
    atoms = render_pptx._tokenize_runs([{"text": "internationalization", "fontSize": 14}],
                                       _content_defaults())
    lines = render_pptx._wrap_lines(atoms, 60, 60)  # 词宽 ~154px > 60px
    assert len(lines) == 1


def test_wrap_trailing_space_ignored():
    base = render_pptx._tokenize_runs([{"text": "字" * 17, "fontSize": 14}], _content_defaults())
    padded = render_pptx._tokenize_runs([{"text": "字" * 17 + "   ", "fontSize": 14}],
                                        _content_defaults())
    assert len(render_pptx._wrap_lines(base, 250, 250)) == \
        len(render_pptx._wrap_lines(padded, 250, 250)) == 1


# --- 行高模型 ---------------------------------------------------------------


def test_paragraph_height_line_factors():
    d = _content_defaults()
    para = [{"runs": [{"text": "你好", "fontSize": 14}]}]
    assert render_pptx.estimate_paragraphs_height(para, d, 250) == pytest.approx(14 * render_pptx.LINE_SINGLE_FACTOR)
    para15 = [{"runs": [{"text": "你好", "fontSize": 14}], "lineHeight": 1.5}]
    assert render_pptx.estimate_paragraphs_height(para15, d, 250) == pytest.approx(14 * render_pptx.LINE_SINGLE_FACTOR * 1.5)


def test_paragraph_space_before_after():
    d = _content_defaults()
    para = [{"runs": [{"text": "你好", "fontSize": 14}], "spaceBefore": 4, "spaceAfter": 6}]
    assert render_pptx.estimate_paragraphs_height(para, d, 250) == pytest.approx(14 * render_pptx.LINE_SINGLE_FACTOR + 10)


def test_bullet_indent_narrows_wrap():
    """bullet 首行少 0.8em、次行起少 18.75px：34 字在 250px 宽，无 bullet 2 行、带 bullet 3 行。"""
    d = _content_defaults()
    plain = [{"runs": [{"text": "字" * 34, "fontSize": 14}]}]
    bullet = [{"runs": [{"text": "字" * 34, "fontSize": 14}], "bullet": True}]
    assert render_pptx.estimate_paragraphs_height(plain, d, 250) == pytest.approx(2 * 14 * render_pptx.LINE_SINGLE_FACTOR)
    assert render_pptx.estimate_paragraphs_height(bullet, d, 250) == pytest.approx(3 * 14 * render_pptx.LINE_SINGLE_FACTOR)


# --- 溢出分级（validate_deck） ----------------------------------------------


def test_compact_single_line_bar_passes():
    """kicker/takeaway 型单行条（height 24 < 文字高+上下inset）静默放行。"""
    el = {"id": "k1", "type": "text", "left": 40, "top": 40, "width": 200, "height": 24,
          "paragraphs": [{"runs": [{"text": "数据工程", "fontSize": 13}]}]}
    assert _codes(_issues_for(el)) == []


def test_overflow_fixable_is_warning():
    el = {"id": "c1", "type": "text", "left": 40, "top": 40, "width": 300, "height": 125,
          "paragraphs": [{"runs": [{"text": "字" * 138, "fontSize": 14}]}]}
    issues = _issues_for(el)
    assert "overflow.fixable" in _codes(issues)
    i = next(i for i in issues if i["code"] == "overflow.fixable")
    assert i["severity"] == "warning"
    assert "自动缩字" in i["message"] and "高度改 ≥" in i["message"]


def test_overflow_unfixable_is_error():
    el = {"id": "c2", "type": "text", "left": 40, "top": 40, "width": 300, "height": 120,
          "paragraphs": [{"runs": [{"text": "字" * 600, "fontSize": 14}]}]}
    issues = _issues_for(el)
    assert "overflow.unfixable" in _codes(issues)
    i = next(i for i in issues if i["code"] == "overflow.unfixable")
    assert i["severity"] == "error"
    assert "删减至约" in i["message"]


def test_overflow_autofit_none_is_error():
    el = {"id": "c3", "type": "text", "left": 40, "top": 40, "width": 300, "height": 125,
          "autoFit": "none",
          "paragraphs": [{"runs": [{"text": "字" * 138, "fontSize": 14}]}]}
    issues = _issues_for(el)
    assert _codes(issues) == ["overflow.nofit"]
    assert issues[0]["severity"] == "error"


def test_small_font_never_shrinks():
    """notes 10px 溢出：字号已低于 12px 下限，不许缩 → unfixable error（改文案而非缩字）。"""
    el = {"id": "n1", "type": "text", "left": 40, "top": 40, "width": 300, "height": 40,
          "textType": "notes",
          "paragraphs": [{"runs": [{"text": "字" * 200, "fontSize": 10}]}]}
    plan = render_pptx.plan_autofit(el["paragraphs"], el, THEME, "notes")
    assert plan["scale_min"] == 1.0 and plan["scale"] is None
    assert "overflow.unfixable" in _codes(_issues_for(el))


def test_shape_text_overflow_checked():
    """shape 内嵌文字同样受检（样式从 text_spec、几何从外层元素读）。"""
    el = {"id": "s1", "type": "shape", "left": 40, "top": 40, "width": 300, "height": 120,
          "shape": "roundRect",
          "text": {"paragraphs": [{"runs": [{"text": "字" * 600, "fontSize": 14}]}]}}
    assert "overflow.unfixable" in _codes(_issues_for(el))


# --- 表格 -------------------------------------------------------------------


def test_table_row_heights():
    el = {"width": 400, "cellMinHeight": 36, "colWidths": [0.5, 0.5], "data": [
        [{"text": "短"}, {"text": "短"}],
        [{"text": "字" * 30}, {"text": "短"}],  # 列宽 200-16=184 → 14 字/行 → 3 行
    ]}
    rh = render_pptx.estimate_table_row_heights(el, THEME)
    assert rh[0] == 36.0  # 短文本不撑行
    assert rh[1] == pytest.approx(3 * 13 * render_pptx.LINE_SINGLE_FACTOR + 8)  # 逐内容行数 × 默认 13px


def test_table_canvas_overflow_is_error():
    rows = [[{"text": f"r{i}c{j}"} for j in range(4)] for i in range(12)]
    el = {"id": "t1", "type": "table", "left": 40, "top": 480, "width": 900, "height": 80,
          "data": rows}
    issues = _issues_for(el)
    assert "overflow.table-canvas" in _codes(issues)
    assert next(i for i in issues if i["code"] == "overflow.table-canvas")["severity"] == "error"


def test_table_grow_past_declared_height_is_warning():
    rows = [[{"text": "字" * 40}] for _ in range(4)]
    el = {"id": "t2", "type": "table", "left": 40, "top": 100, "width": 300, "height": 100,
          "data": rows}
    issues = _issues_for(el)
    assert "overflow.table-grow" in _codes(issues)
    assert next(i for i in issues if i["code"] == "overflow.table-grow")["severity"] == "warning"


# --- 渲染 XML：autofit 写入 ---------------------------------------------------


def _render_text_el(el):
    import pptx
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return render_pptx.render_text(slide, el, THEME, 12192.0)


def _body_pr(shape):
    return shape.text_frame._txBody.bodyPr


def test_normautofit_written_with_explicit_scale():
    """溢出可救 → 显式 fontScale 的 normAutofit，且模板残留的 spAutoFit 已剥。"""
    box = _render_text_el({"id": "c1", "type": "text", "left": 40, "top": 40,
                           "width": 300, "height": 125,
                           "paragraphs": [{"runs": [{"text": "字" * 138, "fontSize": 14}]}]})
    body = _body_pr(box)
    fit = body.find(QN("a:normAutofit"))
    assert fit is not None and fit.get("fontScale") == "86000"  # 86%
    assert body.find(QN("a:spAutoFit")) is None
    assert body.find(QN("a:noAutofit")) is None


def test_noautofit_is_default():
    """不溢出的文本框：固定盒（noAutofit），剥掉 add_textbox 自带的 spAutoFit。"""
    box = _render_text_el({"id": "k1", "type": "text", "left": 40, "top": 40,
                           "width": 300, "height": 60,
                           "paragraphs": [{"runs": [{"text": "你好世界", "fontSize": 14}]}]})
    body = _body_pr(box)
    assert body.find(QN("a:noAutofit")) is not None
    assert body.find(QN("a:spAutoFit")) is None
    assert body.find(QN("a:normAutofit")) is None


def test_autofit_floor_keeps_original_size():
    """10px notes 溢出：低于缩字下限 → 保持原字号（无 normAutofit，--check 报 error）。"""
    box = _render_text_el({"id": "n1", "type": "text", "left": 40, "top": 40,
                           "width": 300, "height": 40, "textType": "notes",
                           "paragraphs": [{"runs": [{"text": "字" * 200, "fontSize": 10}]}]})
    body = _body_pr(box)
    assert body.find(QN("a:normAutofit")) is None


def test_shape_text_autofit_written():
    import pptx
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    el = {"id": "s1", "type": "shape", "left": 40, "top": 40, "width": 300, "height": 125,
          "shape": "roundRect",
          "text": {"paragraphs": [{"runs": [{"text": "字" * 138, "fontSize": 14}]}]}}
    shape = render_pptx.render_shape(slide, el, THEME, 12192.0)
    fit = _body_pr(shape).find(QN("a:normAutofit"))
    assert fit is not None and fit.get("fontScale")


def test_table_row_heights_written():
    import pptx
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    el = {"id": "t1", "type": "table", "left": 40, "top": 60, "width": 400, "height": 100,
          "data": [[{"text": "短"}, {"text": "短"}],
                   [{"text": "字" * 30}, {"text": "短"}]]}
    frame = render_pptx.render_table(slide, el, THEME, 12192.0)
    expected = render_pptx.estimate_table_row_heights(el, THEME)
    for r, row in enumerate(frame.table.rows):
        assert row.height / 12192.0 == pytest.approx(expected[r], abs=0.5)


# --- --check --json ----------------------------------------------------------


def test_check_json_output(tmp_path):
    deck = {"canvas": {"width": 1000, "height": 562.5},
            "slides": [{"id": "s1", "elements": [
                {"id": "c2", "type": "text", "left": 40, "top": 40, "width": 300, "height": 120,
                 "paragraphs": [{"runs": [{"text": "字" * 600, "fontSize": 14}]}]},
            ]}]}
    p = tmp_path / "deck.json"
    p.write_text(json.dumps(deck, ensure_ascii=False), encoding="utf-8")
    r = subprocess.run([sys.executable, str(SCRIPTS / "render_pptx.py"), str(p), "--check", "--json"],
                       capture_output=True, text=True, timeout=120)
    data = json.loads(r.stdout)
    assert r.returncode == 1 and data["ok"] is False
    assert any(i["code"] == "overflow.unfixable" for i in data["issues"])
    assert all(set(i) >= {"severity", "code", "message"} for i in data["issues"])


def test_check_json_ok(tmp_path):
    deck = {"canvas": {"width": 1000, "height": 562.5},
            "slides": [{"id": "s1", "elements": [
                {"id": "k1", "type": "text", "left": 40, "top": 40, "width": 300, "height": 60,
                 "paragraphs": [{"runs": [{"text": "你好世界", "fontSize": 14}]}]},
            ]}]}
    p = tmp_path / "deck.json"
    p.write_text(json.dumps(deck, ensure_ascii=False), encoding="utf-8")
    r = subprocess.run([sys.executable, str(SCRIPTS / "render_pptx.py"), str(p), "--check", "--json"],
                       capture_output=True, text=True, timeout=120)
    data = json.loads(r.stdout)
    assert r.returncode == 0 and data["ok"] is True and data["issues"] == []
